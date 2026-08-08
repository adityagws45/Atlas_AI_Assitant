"""Document parsers — PDF / TXT / Markdown."""

from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass, field

logger = logging.getLogger("atlas.documents.parser")


@dataclass
class PageText:
    page: int
    text: str


@dataclass
class ParseResult:
    pages: list[PageText] = field(default_factory=list)
    full_text: str = ""
    page_count: int = 0
    error: str = ""
    is_likely_scanned: bool = False


def parse_bytes(*, data: bytes, filename: str, mime_type: str = "") -> ParseResult:
    name = (filename or "").lower()
    mime = (mime_type or "").lower()
    if name.endswith(".pdf") or "pdf" in mime:
        return parse_pdf(data)
    if name.endswith(".docx") or "wordprocessingml" in mime:
        return parse_docx(data)
    if name.endswith(".pptx") or "presentationml" in mime:
        return parse_pptx(data)
    if name.endswith((".md", ".markdown")):
        return parse_text(data, markdown=True)
    return parse_text(data, markdown=False)


def parse_docx(data: bytes) -> ParseResult:
    if not data:
        return ParseResult(error="empty")
    try:
        from docx import Document  # python-docx
    except ImportError:
        return ParseResult(error="docx_backend_missing")
    try:
        document = Document(io.BytesIO(data))
        parts = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]
        text = clean_extracted_text("\n\n".join(parts))
    except Exception as exc:  # noqa: BLE001
        logger.warning("event=docx_parse_failed err=%s", type(exc).__name__)
        return ParseResult(error="corrupt_docx")
    if not text.strip():
        return ParseResult(error="empty")
    return ParseResult(pages=[PageText(page=1, text=text)], full_text=text, page_count=1)


def parse_pptx(data: bytes) -> ParseResult:
    """Basic PowerPoint text extraction when python-pptx is available."""
    if not data:
        return ParseResult(error="empty")
    try:
        from pptx import Presentation
    except ImportError:
        return ParseResult(error="pptx_backend_missing")
    try:
        prs = Presentation(io.BytesIO(data))
        pages: list[PageText] = []
        for i, slide in enumerate(prs.slides, start=1):
            bits: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    bits.append(shape.text.strip())
            page_text = clean_extracted_text("\n".join(b for b in bits if b))
            if page_text:
                pages.append(PageText(page=i, text=page_text))
        full = "\n\n".join(f"[Slide {p.page}]\n{p.text}" for p in pages)
    except Exception as exc:  # noqa: BLE001
        logger.warning("event=pptx_parse_failed err=%s", type(exc).__name__)
        return ParseResult(error="corrupt_pptx")
    if not full.strip():
        return ParseResult(error="empty")
    return ParseResult(pages=pages, full_text=full, page_count=len(pages) or 1)


def parse_pdf(data: bytes) -> ParseResult:
    if not data:
        return ParseResult(error="empty")
    try:
        from pypdf import PdfReader
    except ImportError:
        return ParseResult(error="pdf_backend_missing")

    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        logger.warning("event=pdf_corrupt err=%s", type(exc).__name__)
        return ParseResult(error="corrupt_pdf")

    pages: list[PageText] = []
    total_chars = 0
    try:
        n = len(reader.pages)
    except Exception:
        return ParseResult(error="corrupt_pdf")

    for i, page in enumerate(reader.pages, start=1):
        try:
            raw = page.extract_text() or ""
        except Exception:
            raw = ""
        text = _normalize_page(raw)
        total_chars += len(text)
        pages.append(PageText(page=i, text=text))

    full = "\n\n".join(f"[Page {p.page}]\n{p.text}" for p in pages if p.text.strip())
    scanned = n > 0 and total_chars < max(40, n * 15)
    if scanned and not full.strip():
        return ParseResult(
            pages=pages,
            full_text="",
            page_count=n,
            error="scanned_or_empty",
            is_likely_scanned=True,
        )
    return ParseResult(
        pages=pages,
        full_text=full,
        page_count=n,
        is_likely_scanned=scanned,
    )


def parse_text(data: bytes, *, markdown: bool = False) -> ParseResult:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = data.decode("latin-1")
        except Exception:
            return ParseResult(error="undecodable")
    cleaned = clean_extracted_text(text)
    if not cleaned.strip():
        return ParseResult(error="empty")
    # Treat as single logical page for text/md
    pages = [PageText(page=1, text=cleaned)]
    return ParseResult(pages=pages, full_text=cleaned, page_count=1)


def clean_extracted_text(text: str) -> str:
    if not text:
        return ""
    # Drop repeated form-feed / soft hyphens
    text = text.replace("\x0c", "\n").replace("\u00ad", "")
    # Normalize newlines
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    # De-hyphenate line-break wraps: "invest-\nment" → "investment"
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    # Strip common running headers/footers patterns (repeated short lines)
    lines = text.split("\n")
    cleaned_lines: list[str] = []
    for line in lines:
        s = line.strip()
        if re.fullmatch(r"\d{1,4}", s):
            continue
        if re.fullmatch(r"page\s+\d+(\s+of\s+\d+)?", s, re.IGNORECASE):
            continue
        cleaned_lines.append(line.rstrip())
    text = "\n".join(cleaned_lines)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def _normalize_page(raw: str) -> str:
    return clean_extracted_text(raw or "")
